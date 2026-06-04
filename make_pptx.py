#!/usr/bin/env python3
"""Erzeugt aus einer Kennzahl-/Laufzeit-/Kategorie-Auswahl eine PowerPoint-Datei
mit der gerankten Ergebnistabelle der Union- & Quoniam-Fonds.

Nutzt dieselbe Auswahllogik wie fund_metrics.py (Datenquelle: data/funds.json).

Beispiele:
    python3 make_pptx.py --metric treynor --period 3y --provider quoniam \
        --out treynor_quoniam_3j.pptx
    python3 make_pptx.py --metric sharpe --period 5y \
        --category "Global Large-Cap Blend Equity" --top 25 --out sharpe.pptx

Benötigt: python-pptx  (pip install python-pptx)
"""
from __future__ import annotations

import argparse
import sys

import fund_metrics as fm

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("python-pptx fehlt. Installation: pip install python-pptx")

# Farbschema
UNION = RGBColor(0xE3, 0x00, 0x0F)     # Union-Rot
QUONIAM = RGBColor(0x1F, 0x4E, 0x79)   # Dunkelblau
HEADER_BG = RGBColor(0x20, 0x20, 0x20)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF2, 0xF2, 0xF2)

ROWS_PER_SLIDE = 18


def _provider_color(branding: str) -> RGBColor:
    return QUONIAM if branding.lower().startswith("quoniam") else UNION


def build_presentation(data: dict, metric: str, period: str,
                       category: str | None, provider: str, top: int | None):
    rows = fm.query(data["funds"], metric, period, category, provider)
    if top:
        rows = rows[:top]
    label = fm.METRICS[metric][0]
    period_label = fm.PERIOD_LABELS[period]
    meta = data.get("meta", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    # --- Titelfolie ---
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text(s, f"{label} – {period_label}", Inches(0.6), Inches(0.7),
              Inches(12), Inches(1.1), size=40, bold=True)
    prov_txt = {"all": "Union Investment + Quoniam", "union": "Union Investment",
                "quoniam": "Quoniam"}[provider]
    cat_txt = category or "Alle Kategorien"
    sub = (f"Anbieter: {prov_txt}    |    Kategorie: {cat_txt}\n"
           f"{len(rows)} Fonds    |    Datenquelle: Morningstar    |    "
           f"Stand: {meta.get('as_of', '?')}")
    _add_text(s, sub, Inches(0.6), Inches(1.9), Inches(12), Inches(1.2),
              size=16, color=RGBColor(0x55, 0x55, 0x55))
    _add_text(s, "Höhere Werte = besser. Ranking absteigend.",
              Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
              size=12, color=RGBColor(0x88, 0x88, 0x88))

    if not rows:
        _add_text(s, "Keine Fonds mit Werten für diese Auswahl.",
                  Inches(0.6), Inches(3.2), Inches(12), Inches(1), size=22, bold=True)
        return prs, rows

    # --- Tabellenfolien (paginiert) ---
    col_w = [Inches(0.9), Inches(7.4), Inches(2.8), Inches(2.0)]
    headers = ["Rang", "Fonds", "Anbieter", f"{label} ({period_label})"]
    for start in range(0, len(rows), ROWS_PER_SLIDE):
        chunk = rows[start:start + ROWS_PER_SLIDE]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text(slide, f"{label} – {period_label}", Inches(0.5), Inches(0.25),
                  Inches(12.3), Inches(0.6), size=22, bold=True)
        tbl = slide.shapes.add_table(
            len(chunk) + 1, 4, Inches(0.5), Inches(1.0),
            sum(col_w, Inches(0)), Inches(0.4) * (len(chunk) + 1)).table
        for i, w in enumerate(col_w):
            tbl.columns[i].width = w
        # Kopfzeile
        for c, h in enumerate(headers):
            _style_cell(tbl.cell(0, c), h, bold=True, fg=HEADER_FG, bg=HEADER_BG,
                        align=(PP_ALIGN.RIGHT if c == 3 else PP_ALIGN.LEFT))
        # Datenzeilen
        for r, row in enumerate(chunk, start=1):
            bg = ROW_ALT if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            _style_cell(tbl.cell(r, 0), str(start + r), bg=bg, align=PP_ALIGN.CENTER)
            _style_cell(tbl.cell(r, 1), row["name"], bg=bg)
            _style_cell(tbl.cell(r, 2), row["branding"], bg=bg,
                        fg=_provider_color(row["branding"]), bold=True)
            _style_cell(tbl.cell(r, 3), f"{row['value']:.3f}", bg=bg,
                        align=PP_ALIGN.RIGHT)
    return prs, rows


def _add_text(slide, text, left, top, width, height, size=18, bold=False,
              color=RGBColor(0, 0, 0)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for j, line in enumerate(text.split("\n")):
        para = p if j == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _style_cell(cell, text, bold=False, fg=RGBColor(0x20, 0x20, 0x20),
                bg=None, align=PP_ALIGN.LEFT):
    cell.text = text
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    para = cell.text_frame.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(12)
    run.font.bold = bold
    run.font.color.rgb = fg
    cell.margin_top = Pt(1)
    cell.margin_bottom = Pt(1)


def main(argv=None):
    p = argparse.ArgumentParser(description="PowerPoint-Export der Fonds-Kennzahlen.")
    p.add_argument("--metric", required=True, choices=list(fm.METRICS.keys()))
    p.add_argument("--period", required=True, choices=list(fm.PERIOD_LABELS.keys()))
    p.add_argument("--category", help="Morningstar-Kategorie (ohne 'EAA Fund'-Präfix)")
    p.add_argument("--provider", choices=["all", "union", "quoniam"], default="all")
    p.add_argument("--top", type=int, default=None)
    p.add_argument("--out", required=True, help="Ausgabedatei (.pptx)")
    p.add_argument("--data", default=fm.DATA_FILE)
    args = p.parse_args(argv)

    data = fm.load_data(args.data)
    try:
        prs, rows = build_presentation(data, args.metric, args.period,
                                       args.category, args.provider, args.top)
    except ValueError as exc:
        sys.exit(f"Fehler: {exc}")
    prs.save(args.out)
    print(f"{len(rows)} Fonds -> {args.out}")


if __name__ == "__main__":
    main()
